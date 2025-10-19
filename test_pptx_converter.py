"""
Демонстрационный скрипт для создания тестовой PPTX презентации
и последующей конвертации в HTML
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def create_demo_presentation():
    """Создаёт демонстрационную PPTX презентацию"""
    
    print("🎨 Создание демонстрационной презентации...")
    
    # Создаём презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # === СЛАЙД 1: Титульный ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    
    # Фон
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(50, 56, 92)  # Тёмно-синий
    
    # Заголовок
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "PPTX to HTML Converter"
    
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 252, 247)  # Кремовый
    
    # Подзаголовок
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Демонстрация возможностей конвертации"
    
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(129, 134, 154)  # Серо-голубой
    
    print("  ✅ Слайд 1: Титульный")
    
    # === СЛАЙД 2: Текстовое форматирование ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Фон
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 252, 247)
    
    # Заголовок
    heading_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    heading_frame = heading_box.text_frame
    heading_frame.text = "Текстовое форматирование"
    heading_para = heading_frame.paragraphs[0]
    heading_para.font.size = Pt(40)
    heading_para.font.bold = True
    heading_para.font.color.rgb = RGBColor(50, 56, 92)
    
    # Текстовый блок с разным форматированием
    text_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    )
    tf = text_box.text_frame
    
    # Первый параграф
    p1 = tf.paragraphs[0]
    p1.text = "Обычный текст с разными стилями:"
    p1.font.size = Pt(24)
    p1.font.color.rgb = RGBColor(0, 0, 0)
    
    # Второй параграф - жирный
    p2 = tf.add_paragraph()
    p2.text = "Жирный текст"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 0, 0)
    p2.level = 1
    
    # Третий параграф - курсив
    p3 = tf.add_paragraph()
    p3.text = "Курсивный текст"
    p3.font.size = Pt(20)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0, 128, 0)
    p3.level = 1
    
    # Четвёртый параграф - подчёркнутый
    p4 = tf.add_paragraph()
    p4.text = "Подчёркнутый текст"
    p4.font.size = Pt(20)
    p4.font.underline = True
    p4.font.color.rgb = RGBColor(0, 0, 255)
    p4.level = 1
    
    # Пятый параграф - комбинация
    p5 = tf.add_paragraph()
    p5.text = "Жирный курсивный подчёркнутый"
    p5.font.size = Pt(20)
    p5.font.bold = True
    p5.font.italic = True
    p5.font.underline = True
    p5.font.color.rgb = RGBColor(128, 0, 128)
    p5.level = 1
    
    # Выравнивание
    p6 = tf.add_paragraph()
    p6.text = "Текст по центру"
    p6.font.size = Pt(18)
    p6.alignment = PP_ALIGN.CENTER
    
    p7 = tf.add_paragraph()
    p7.text = "Текст справа"
    p7.font.size = Pt(18)
    p7.alignment = PP_ALIGN.RIGHT
    
    print("  ✅ Слайд 2: Текстовое форматирование")
    
    # === СЛАЙД 3: Фигуры и цвета ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Заголовок
    heading3 = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    heading3.text_frame.text = "Фигуры и цветные блоки"
    heading3.text_frame.paragraphs[0].font.size = Pt(40)
    heading3.text_frame.paragraphs[0].font.bold = True
    heading3.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 56, 92)
    
    # Цветные блоки
    colors = [
        (RGBColor(255, 0, 0), "Красный"),
        (RGBColor(0, 255, 0), "Зелёный"),
        (RGBColor(0, 0, 255), "Синий"),
        (RGBColor(255, 255, 0), "Жёлтый"),
    ]
    
    for i, (color, name) in enumerate(colors):
        # Цветной блок
        shape = slide3.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5 + i * 2.3), Inches(2),
            Inches(2), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        # Текст на блоке
        text_frame = shape.text_frame
        text_frame.text = name
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("  ✅ Слайд 3: Фигуры и цвета")
    
    # === СЛАЙД 4: Списки ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Заголовок
    heading4 = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    heading4.text_frame.text = "Списки и уровни"
    heading4.text_frame.paragraphs[0].font.size = Pt(40)
    heading4.text_frame.paragraphs[0].font.bold = True
    
    # Список
    list_box = slide4.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(5)
    )
    tf4 = list_box.text_frame
    
    items = [
        (0, "Первый уровень - пункт 1"),
        (1, "Второй уровень - подпункт 1.1"),
        (1, "Второй уровень - подпункт 1.2"),
        (2, "Третий уровень - подпункт 1.2.1"),
        (0, "Первый уровень - пункт 2"),
        (1, "Второй уровень - подпункт 2.1"),
    ]
    
    tf4.paragraphs[0].text = items[0][1]
    tf4.paragraphs[0].level = items[0][0]
    tf4.paragraphs[0].font.size = Pt(20)
    
    for level, text in items[1:]:
        p = tf4.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20)
    
    print("  ✅ Слайд 4: Списки")
    
    # === СЛАЙД 5: Финальный ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Фон
    background5 = slide5.background
    fill5 = background5.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(50, 56, 92)
    
    # Текст
    final_box = slide5.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(2)
    )
    final_frame = final_box.text_frame
    final_frame.text = "Спасибо!\n\nВсе элементы успешно конвертированы"
    
    for para in final_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 252, 247)
    
    print("  ✅ Слайд 5: Финальный")
    
    # Сохраняем презентацию
    filename = "demo_presentation.pptx"
    prs.save(filename)
    
    print(f"\n✅ Презентация создана: {filename}")
    return filename


def test_converter():
    """Тестирует конвертер на демо-презентации"""
    print("\n" + "=" * 60)
    print("ТЕСТИРОВАНИЕ PPTX TO HTML CONVERTER")
    print("=" * 60 + "\n")
    
    # Создаём демо-презентацию
    pptx_file = create_demo_presentation()
    
    # Проверяем наличие конвертера
    if not os.path.exists('pptx_to_html.py'):
        print("❌ Файл pptx_to_html.py не найден!")
        return
    
    print("\n" + "=" * 60)
    print("КОНВЕРТАЦИЯ В HTML")
    print("=" * 60 + "\n")
    
    # Импортируем и запускаем конвертер
    from pptx_to_html import PPTXToHTMLConverter
    
    converter = PPTXToHTMLConverter(pptx_file, 'demo_output')
    converter.convert()
    
    print("\n" + "=" * 60)
    print("✨ ТЕСТ ЗАВЕРШЁН УСПЕШНО!")
    print("=" * 60)
    print("\n📂 Результаты:")
    print(f"   - Исходный PPTX: {pptx_file}")
    print(f"   - HTML вывод: demo_output/index.html")
    print(f"   - CSS файл: demo_output/style.css")
    print(f"   - Изображения: demo_output/images/")
    print("\n🌐 Откройте demo_output/index.html в браузере!")
    print()


if __name__ == '__main__':
    test_converter()
