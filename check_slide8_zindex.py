#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Проверка z-index на слайде 8 с текущей логикой
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_shape_info(shape, slide_width, slide_height):
    """Получить информацию о фигуре"""
    left_pct = (shape.left / slide_width) * 100
    top_pct = (shape.top / slide_height) * 100
    width_pct = (shape.width / slide_width) * 100
    height_pct = (shape.height / slide_height) * 100
    
    shape_type = MSO_SHAPE_TYPE(shape.shape_type).name
    
    is_text = False
    text_preview = ""
    if hasattr(shape, 'text'):
        try:
            if shape.text.strip():
                is_text = True
                text_preview = shape.text[:40] + "..." if len(shape.text) > 40 else shape.text
        except:
            pass
    
    return {
        'name': shape.name,
        'type': shape_type,
        'left': left_pct,
        'top': top_pct,
        'width': width_pct,
        'height': height_pct,
        'right': left_pct + width_pct,
        'bottom': top_pct + height_pct,
        'is_text': is_text,
        'text': text_preview
    }

# Загружаем презентацию
prs = Presentation('NEO INVESTMENTS-fin.pptx')
slide = prs.slides[7]  # Слайд 8

slide_width = prs.slide_width
slide_height = prs.slide_height

print("=" * 70)
print("ПРОВЕРКА Z-INDEX НА СЛАЙДЕ 8")
print("=" * 70)
print("\nТекущая логика: z-index = shape_counter (порядок обработки)\n")

shape_counter = 0
all_blocks = []

def process_shape(shape, level=0):
    global shape_counter
    indent = "  " * level
    
    # Пропускаем GROUP сам по себе, обрабатываем содержимое
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}GROUP: {shape.name}")
        for sub_shape in shape.shapes:
            process_shape(sub_shape, level + 1)
        return
    
    shape_counter += 1
    info = get_shape_info(shape, slide_width, slide_height)
    info['z_index'] = shape_counter
    all_blocks.append(info)
    
    # Определяем положение
    pos = []
    if info['left'] < 30:
        pos.append("ЛЕВЫЙ")
    elif info['left'] > 60:
        pos.append("ПРАВЫЙ")
    
    if info['top'] < 30:
        pos.append("ВЕРХНИЙ")
    elif info['top'] > 60:
        pos.append("НИЖНИЙ")
    else:
        pos.append("СРЕДНИЙ")
    
    pos_str = " ".join(pos) if pos else "ЦЕНТР"
    
    type_str = "ТЕКСТ" if info['is_text'] else "ИЗОБРАЖЕНИЕ"
    
    print(f"{indent}#{shape_counter}: {info['name']} - {type_str} - {pos_str}")
    print(f"{indent}        Position: left={info['left']:.1f}%, top={info['top']:.1f}%")
    print(f"{indent}        z-index: {shape_counter}")
    if info['is_text']:
        print(f"{indent}        Text: {info['text']}")

print()
for shape in slide.shapes:
    process_shape(shape)

print("\n" + "=" * 70)
print("АНАЛИЗ ПЕРЕСЕЧЕНИЙ:")
print("=" * 70)

# Проверяем пересечения
images = [b for b in all_blocks if not b['is_text']]

for i in range(len(images)):
    for j in range(i + 1, len(images)):
        b1 = images[i]
        b2 = images[j]
        
        # Проверяем пересечение
        overlap_left = max(b1['left'], b2['left'])
        overlap_right = min(b1['right'], b2['right'])
        overlap_top = max(b1['top'], b2['top'])
        overlap_bottom = min(b1['bottom'], b2['bottom'])
        
        overlap_width = max(0, overlap_right - overlap_left)
        overlap_height = max(0, overlap_bottom - overlap_top)
        
        if overlap_width > 1 and overlap_height > 1:  # Минимум 1% пересечения
            print(f"\n⚠️ ПЕРЕСЕЧЕНИЕ:")
            print(f"   Блок '{b1['name']}' (z={b1['z_index']}, top={b1['top']:.1f}%)")
            print(f"   Блок '{b2['name']}' (z={b2['z_index']}, top={b2['top']:.1f}%)")
            
            if b1['z_index'] < b2['z_index']:
                print(f"   ✅ '{b2['name']}' ПОВЕРХ '{b1['name']}'")
            else:
                print(f"   ❌ ПРОБЛЕМА: '{b1['name']}' поверх '{b2['name']}'")

print("\n" + "=" * 70)
