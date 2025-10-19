"""
Поиск фигур с границами и детальный анализ их стилей
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.util import Pt

def rgb_to_hex(color_obj):
    """Конвертирует цвет в hex"""
    try:
        if color_obj.type == MSO_COLOR_TYPE.RGB:
            rgb = color_obj.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        elif color_obj.type == MSO_COLOR_TYPE.SCHEME:
            try:
                rgb = color_obj.rgb
                return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
            except:
                return None
    except:
        return None

def emu_to_px(emu):
    return emu / 9525

# Загружаем презентацию
prs = Presentation("NEO INVESTMENTS-fin.pptx")

print("="*70)
print("ПОИСК ФИГУР С ГРАНИЦАМИ И ДРУГИМИ СТИЛЯМИ")
print("="*70)

# Проверяем слайд 3 (у него много фигур с границами)
for slide_num in [3, 4, 6]:
    slide = prs.slides[slide_num - 1]
    
    print(f"\n{'='*70}")
    print(f"СЛАЙД {slide_num} - {len(slide.shapes)} фигур")
    print(f"{'='*70}")
    
    found_borders = []
    found_transparency = []
    found_no_fill = []
    
    for idx, shape in enumerate(slide.shapes, 1):
        try:
            shape_info = {
                'idx': idx,
                'type': shape.shape_type.name if hasattr(shape, 'shape_type') else 'Unknown',
                'name': shape.name if hasattr(shape, 'name') else 'N/A',
                'has_border': False,
                'border_info': None,
                'has_fill': False,
                'fill_info': None,
                'transparency': None
            }
            
            # Проверяем границу
            if hasattr(shape, 'line'):
                line = shape.line
                
                # Проверяем цвет границы
                if hasattr(line, 'color') and line.color:
                    try:
                        border_color = rgb_to_hex(line.color)
                        if border_color:
                            shape_info['has_border'] = True
                            
                            # Толщина
                            width = "None"
                            if hasattr(line, 'width') and line.width:
                                try:
                                    width_pt = line.width.pt
                                    width = f"{width_pt:.1f}pt"
                                except:
                                    width_emu = line.width
                                    width = f"{emu_to_px(width_emu):.2f}px"
                            
                            shape_info['border_info'] = f"{border_color} {width}"
                            found_borders.append(shape_info)
                    except:
                        pass
            
            # Проверяем заливку
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill.type == MSO_FILL_TYPE.SOLID:
                    shape_info['has_fill'] = True
                    fill_color = rgb_to_hex(fill.fore_color)
                    shape_info['fill_info'] = f"SOLID {fill_color}"
                    
                    # Проверяем прозрачность
                    try:
                        if hasattr(fill.fore_color, 'transparency'):
                            trans = fill.fore_color.transparency
                            if trans is not None and trans > 0:
                                shape_info['transparency'] = trans
                                found_transparency.append(shape_info)
                    except:
                        pass
                        
                elif fill.type is None or fill.type == MSO_FILL_TYPE.BACKGROUND:
                    found_no_fill.append(shape_info)
                    
        except Exception as e:
            continue
    
    # Выводим результаты
    if found_borders:
        print(f"\n  Фигуры с границами ({len(found_borders)}):")
        for s in found_borders[:5]:  # Первые 5
            print(f"    #{s['idx']} ({s['type']}): border = {s['border_info']}")
            if s['fill_info']:
                print(f"             fill = {s['fill_info']}")
    
    if found_transparency:
        print(f"\n  Фигуры с прозрачностью ({len(found_transparency)}):")
        for s in found_transparency[:3]:
            print(f"    #{s['idx']} ({s['type']}): transparency = {s['transparency']:.2f}")
    
    if found_no_fill:
        print(f"\n  Фигуры без заливки ({len(found_no_fill)}):")
        for s in found_no_fill[:3]:
            print(f"    #{s['idx']} ({s['type']} '{s['name']}')")

print(f"\n{'='*70}")
print("ИТОГ: Найдены фигуры с границами, прозрачностью и без заливки")
print("Эти стили должны быть извлечены в extract_shape_style()")
print("="*70)
