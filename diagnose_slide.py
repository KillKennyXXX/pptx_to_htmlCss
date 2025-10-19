"""
Диагностический скрипт для анализа слайда
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

def diagnose_slide(pptx_path, slide_num):
    """Диагностика конкретного слайда"""
    
    print(f"\n{'='*60}")
    print(f"ДИАГНОСТИКА СЛАЙДА {slide_num}")
    print(f"{'='*60}\n")
    
    prs = Presentation(pptx_path)
    
    if slide_num > len(prs.slides):
        print(f"Ошибка: слайд {slide_num} не существует (всего {len(prs.slides)} слайдов)")
        return
    
    slide = prs.slides[slide_num - 1]
    
    # Размеры слайда
    slide_width = prs.slide_width / 914400  # EMU to inches
    slide_height = prs.slide_height / 914400
    print(f"📏 Размеры слайда: {slide_width:.2f}\" x {slide_height:.2f}\"")
    print(f"   В пикселях: {int(prs.slide_width/12700)} x {int(prs.slide_height/12700)} px\n")
    
    # Фон
    print("🎨 ФОН СЛАЙДА:")
    try:
        if hasattr(slide.background, 'fill'):
            fill = slide.background.fill
            print(f"   Тип заливки: {fill.type} ({type(fill.type)})")
            
            if fill.type == 1:  # SOLID
                print("   → Сплошной цвет")
            elif fill.type == 6:  # PICTURE
                print("   → ИЗОБРАЖЕНИЕ!")
                # Пытаемся найти изображение в XML
                bg_element = slide.background._element
                print(f"   XML элемент: {bg_element.tag}")
                
                # Проверяем relationships
                print(f"\n   Relationships слайда:")
                for rel_id, rel in slide.part.rels.items():
                    if 'image' in rel.reltype.lower():
                        print(f"      {rel_id}: {rel.target_ref} ({rel.reltype})")
            else:
                print(f"   → Другой тип: {fill.type}")
        else:
            print("   Нет заливки")
    except Exception as e:
        print(f"   ⚠️ Ошибка: {e}")
    
    # Список всех фигур
    print(f"\n📦 ФИГУРЫ НА СЛАЙДЕ (всего: {len(slide.shapes)}):\n")
    
    shape_types_count = {}
    
    def analyze_shape(shape, level=0):
        """Рекурсивный анализ фигуры"""
        indent = "  " * level
        
        # Тип фигуры
        shape_type = shape.shape_type
        type_name = str(shape_type).split('.')[-1].replace('(', '').replace(')', '')
        
        if shape_type not in shape_types_count:
            shape_types_count[shape_type] = 0
        shape_types_count[shape_type] += 1
        
        # Позиция и размер
        try:
            left = shape.left / 914400  # EMU to inches
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            pos_info = f"[{left:.1f}\", {top:.1f}\"] {width:.1f}\"x{height:.1f}\""
        except:
            pos_info = "[нет позиции]"
        
        # Текст
        text_info = ""
        if hasattr(shape, 'text') and shape.text:
            text_preview = shape.text[:30].replace('\n', ' ')
            if len(shape.text) > 30:
                text_preview += "..."
            text_info = f" | Текст: '{text_preview}'"
        
        # Имя фигуры
        name_info = ""
        if hasattr(shape, 'name'):
            name_info = f" | Имя: {shape.name}"
        
        print(f"{indent}• {type_name:20} {pos_info:25}{name_info}{text_info}")
        
        # Если это картинка
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                print(f"{indent}  └─ Изображение: {img.ext} ({len(img.blob)} байт)")
            except Exception as e:
                print(f"{indent}  └─ ⚠️ Ошибка получения изображения: {e}")
        
        # Если это группа
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{indent}  └─ Группа содержит {len(shape.shapes)} фигур:")
            for sub_shape in shape.shapes:
                analyze_shape(sub_shape, level + 2)
        
        # Если это FREEFORM - проверяем заливку
        if shape_type == MSO_SHAPE_TYPE.FREEFORM:
            try:
                if hasattr(shape, 'fill'):
                    fill_type = shape.fill.type
                    print(f"{indent}  └─ Заливка: тип {fill_type} ({type(fill_type)})")
                    if fill_type == 6:  # PICTURE
                        print(f"{indent}      ⚠️ FREEFORM с заливкой-изображением!")
            except Exception as e:
                print(f"{indent}  └─ Ошибка проверки заливки: {e}")
    
    for i, shape in enumerate(slide.shapes, 1):
        print(f"\n[{i}]")
        analyze_shape(shape, 1)
    
    # Статистика
    print(f"\n\n📊 СТАТИСТИКА:")
    print(f"   Всего фигур: {len(slide.shapes)}")
    print(f"\n   По типам:")
    for shape_type, count in sorted(shape_types_count.items(), key=lambda x: x[1], reverse=True):
        type_name = str(shape_type).split('.')[-1].replace('(', '').replace(')', '')
        print(f"      {type_name:20} x {count}")
    
    # Ищем самые большие изображения
    print(f"\n\n🔍 АНАЛИЗ ИЗОБРАЖЕНИЙ:")
    images = []
    
    def collect_images(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                area = shape.width * shape.height
                images.append({
                    'shape': shape,
                    'area': area,
                    'width': shape.width / 914400,
                    'height': shape.height / 914400
                })
            except:
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                collect_images(sub)
    
    for shape in slide.shapes:
        collect_images(shape)
    
    if images:
        slide_area = prs.slide_width * prs.slide_height
        images_sorted = sorted(images, key=lambda x: x['area'], reverse=True)
        
        print(f"   Найдено изображений: {len(images)}")
        print(f"\n   Самые большие:")
        for i, img in enumerate(images_sorted[:5], 1):
            percent = (img['area'] / slide_area) * 100
            print(f"      {i}. {img['width']:.1f}\" x {img['height']:.1f}\" (занимает {percent:.1f}% слайда)")
    else:
        print("   Изображения не найдены")
    
    print(f"\n{'='*60}\n")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Использование: python diagnose_slide.py <файл.pptx> <номер_слайда>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])
    
    diagnose_slide(pptx_path, slide_num)
