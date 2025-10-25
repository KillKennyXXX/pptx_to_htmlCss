"""
PPTX to HTML Converter (v17.0)
Конвертирует презентации PowerPoint в веб-страницы с сохранением форматирования

Версия 15: Улучшенная классификация изображений (QR-коды, иконки, логотипы)
Версия 16: Полное извлечение стилей (градиенты, тени, границы, трансформации)
Версия 16.1: Исправлена логика границ и теней (удаление ложных границ)
Версия 16.2: Исправлена прозрачность PNG изображений
Версия 16.3: Добавлена поддержка композитных QR-кодов из групп фигур
Версия 17.0: Каждый слайд сохраняется в отдельный HTML файл (папка pages/)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.dml.color import RGBColor
import os
import base64
from pathlib import Path
import json
import re

# Импортируем классификатор изображений
from image_classifier import ImageClassifier

# v16: Импортируем извлекатель продвинутых стилей
from style_extractor import style_extractor


class PPTXToHTMLConverter:
    def __init__(self, pptx_path, output_dir='pptx_output'):
        """
        Инициализация конвертера
        
        Args:
            pptx_path: Путь к PPTX файлу
            output_dir: Папка для сохранения HTML и изображений
        """
        self.pptx_path = pptx_path
        self.output_dir = output_dir
        self.images_dir = os.path.join(output_dir, 'images')
        self.pages_dir = os.path.join(output_dir, 'pages')
        self.prs = None
        self.slide_data = []
        self.current_slide_bg_color = None  # Для определения дефолтного цвета текста
        
        # v15: Инициализируем классификатор изображений
        self.image_classifier = ImageClassifier()
        
        # Создаем директории
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.images_dir, exist_ok=True)
        os.makedirs(self.pages_dir, exist_ok=True)
    
    def load_presentation(self):
        """Загружает презентацию"""
        print(f"Загрузка презентации: {self.pptx_path}")
        self.prs = Presentation(self.pptx_path)
        print(f"Найдено слайдов: {len(self.prs.slides)}")
    
    def get_default_text_color(self):
        """Определяет дефолтный цвет текста на основе яркости фона слайда"""
        if not self.current_slide_bg_color:
            return '#000000'  # Черный по умолчанию
        
        try:
            # Убираем #
            hex_color = self.current_slide_bg_color.lstrip('#')
            
            # Конвертируем в RGB
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            
            # Вычисляем яркость (luminance) по формуле
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            
            # Если фон темный (luminance < 0.5), используем белый текст
            # Если фон светлый, используем черный текст
            return '#ffffff' if luminance < 0.5 else '#000000'
        except:
            return '#000000'
    
    def rgb_to_hex(self, rgb_color):
        """Конвертирует RGB в HEX с правильной обработкой theme colors"""
        if rgb_color is None:
            return None
        try:
            # Проверяем, что это ColorFormat
            if hasattr(rgb_color, 'type'):
                color_type = rgb_color.type
                
                # Если это RGB color (type = 1) - прямой RGB
                if color_type == 1:  # MSO_COLOR_TYPE.RGB
                    try:
                        if hasattr(rgb_color, 'rgb'):
                            r, g, b = rgb_color.rgb
                            return f'#{r:02x}{g:02x}{b:02x}'
                    except:
                        pass
                
                # Если это theme color (SCHEME) - извлекаем реальный RGB
                elif color_type == 2:  # MSO_COLOR_TYPE.SCHEME
                    try:
                        # Получаем реальный RGB через встроенный метод
                        # Этот метод учитывает brightness, tint и shade
                        if hasattr(rgb_color, 'rgb'):
                            r, g, b = rgb_color.rgb
                            return f'#{r:02x}{g:02x}{b:02x}'
                    except:
                        pass
            
            # Если это прямой RGBColor объект
            if hasattr(rgb_color, 'rgb'):
                r, g, b = rgb_color.rgb
                return f'#{r:02x}{g:02x}{b:02x}'
            
            return None
        except Exception as e:
            return None
    
    def pt_to_px(self, pt_value):
        """Конвертирует points в pixels"""
        if pt_value is None:
            return None
        return round(pt_value * 1.333333)
    
    def emu_to_px(self, emu_value):
        """Конвертирует EMU (English Metric Units) в pixels"""
        if emu_value is None:
            return None
        return round(emu_value / 9525)
    
    def get_text_alignment(self, alignment):
        """Получает CSS выравнивание текста"""
        alignment_map = {
            PP_ALIGN.LEFT: 'left',
            PP_ALIGN.CENTER: 'center',
            PP_ALIGN.RIGHT: 'right',
            PP_ALIGN.JUSTIFY: 'justify',
        }
        return alignment_map.get(alignment, 'left')
    
    def get_vertical_alignment(self, vertical_anchor):
        """Получает CSS вертикальное выравнивание"""
        anchor_map = {
            MSO_VERTICAL_ANCHOR.TOP: 'flex-start',
            MSO_VERTICAL_ANCHOR.MIDDLE: 'center',
            MSO_VERTICAL_ANCHOR.BOTTOM: 'flex-end',
        }
        return anchor_map.get(vertical_anchor, 'flex-start')
    
    def extract_text_formatting(self, run):
        """Извлекает форматирование текста"""
        style = {}
        
        try:
            # Шрифт
            if run.font.name:
                style['font-family'] = run.font.name
            
            # Размер шрифта
            if run.font.size:
                style['font-size'] = f"{self.pt_to_px(run.font.size.pt)}px"
            
            # Цвет текста
            if run.font.color:
                if run.font.color.type:
                    # Явно заданный цвет
                    color = self.rgb_to_hex(run.font.color)
                    if color:
                        style['color'] = color
                else:
                    # Цвет не задан явно - используем дефолтный на основе фона слайда
                    style['color'] = self.get_default_text_color()
            else:
                # Если color вообще None, используем дефолтный на основе фона
                style['color'] = self.get_default_text_color()
            
            # Жирный
            if run.font.bold:
                style['font-weight'] = 'bold'
            
            # Курсив
            if run.font.italic:
                style['font-style'] = 'italic'
            
            # Подчеркивание
            if run.font.underline:
                style['text-decoration'] = 'underline'
        except:
            pass
        
        return style
    
    def extract_paragraph_formatting(self, paragraph):
        """Извлекает форматирование параграфа"""
        style = {}
        
        try:
            # Выравнивание
            if paragraph.alignment:
                style['text-align'] = self.get_text_alignment(paragraph.alignment)
            
            # Отступы
            if paragraph.level > 0:
                style['margin-left'] = f"{paragraph.level * 20}px"
            
            # Межстрочный интервал
            if paragraph.line_spacing:
                style['line-height'] = str(paragraph.line_spacing)
            
            # Отступы до и после параграфа
            if paragraph.space_before:
                style['margin-top'] = f"{self.pt_to_px(paragraph.space_before.pt)}px"
            if paragraph.space_after:
                style['margin-bottom'] = f"{self.pt_to_px(paragraph.space_after.pt)}px"
        except:
            pass
        
        return style
    
    def extract_shape_style(self, shape, slide_width, slide_height, shape_index=0):
        """Извлекает стили формы с процентными размерами для адаптивности
        
        Args:
            shape: Фигура для извлечения стилей
            slide_width: Ширина слайда в пикселях
            slide_height: Высота слайда в пикселях
            shape_index: Индекс фигуры для вычисления z-index (по умолчанию 0)
        
        Note:
            Координаты shape.left и shape.top уже абсолютные относительно слайда,
            даже для фигур внутри групп (python-pptx автоматически это учитывает)
        """
        # Абсолютные значения в пикселях (уже учитывают положение в группе)
        left_px = self.emu_to_px(shape.left)
        top_px = self.emu_to_px(shape.top)
        width_px = self.emu_to_px(shape.width)
        height_px = self.emu_to_px(shape.height)
        
        # Конвертируем в проценты для адаптивности
        left_percent = (left_px / slide_width) * 100
        top_percent = (top_px / slide_height) * 100
        width_percent = (width_px / slide_width) * 100
        height_percent = (height_px / slide_height) * 100
        
        # Вычисляем z-index на основе порядка обработки фигур
        # В PowerPoint порядок фигур в slide.shapes определяет z-order:
        # - Первая фигура (shapes[0]) - самая задняя (z-index минимальный)
        # - Последняя фигура (shapes[-1]) - самая передняя (z-index максимальный)
        # 
        # shape_index соответствует порядку обработки фигур, который совпадает
        # с порядком в slide.shapes (с учетом рекурсивной обработки групп)
        #
        # Используем простую формулу: z-index = shape_index
        # Это сохраняет оригинальный порядок наложения из PowerPoint
        
        z_index = shape_index
        
        style = {
            'position': 'absolute',
            'left': f"{left_percent:.3f}%",
            'top': f"{top_percent:.3f}%",
            'width': f"{width_percent:.3f}%",
            'height': f"{height_percent:.3f}%",
            'z-index': str(z_index),
        }
        
        try:
            # v16: Используем StyleExtractor для продвинутых стилей
            
            # FILL (заливка) - поддержка SOLID, GRADIENT, PATTERN, PICTURE
            if hasattr(shape, 'fill'):
                fill_styles = style_extractor.extract_fill_style(shape.fill)
                style.update(fill_styles)
            
            # LINE (граница) - поддержка разных стилей линий
            if hasattr(shape, 'line'):
                line_styles = style_extractor.extract_line_style(shape.line)
                style.update(line_styles)
            
            # EFFECTS (эффекты) - тени, свечение и т.д.
            shadow_styles = style_extractor.extract_shadow_effect(shape)
            style.update(shadow_styles)
            
            # TRANSFORMS (трансформации) - rotation, flip
            transform_styles = style_extractor.extract_transform_style(shape)
            style.update(transform_styles)
                    
        except Exception as e:
            # Игнорируем ошибки извлечения стилей
            print(f"      ⚠️ Ошибка извлечения стилей: {e}")
            pass
        
        return style
    
    def save_image(self, image, slide_num, img_num, prefix="img"):
        """Сохраняет изображение"""
        try:
            img_filename = f"slide{slide_num}_{prefix}{img_num}.{image.ext}"
            img_path = os.path.join(self.images_dir, img_filename)
            
            with open(img_path, 'wb') as f:
                f.write(image.blob)
            
            return f"images/{img_filename}"
        except Exception as e:
            print(f"Ошибка сохранения изображения: {e}")
            return None
    
    def save_background_image(self, slide, slide_num):
        """Сохраняет фоновое изображение слайда"""
        try:
            # Метод 1: Поиск самого большого изображения (часто фон - это просто большая картинка)
            # Ищем изображения рекурсивно, включая группы
            def find_large_images(shapes_list):
                images = []
                for shp in shapes_list:
                    try:
                        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            width = self.emu_to_px(shp.width)
                            height = self.emu_to_px(shp.height)
                            left = self.emu_to_px(shp.left)
                            top = self.emu_to_px(shp.top)
                            area = width * height
                            images.append({'shape': shp, 'area': area, 'left': left, 'top': top})
                        elif shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                            images.extend(find_large_images(shp.shapes))
                    except:
                        continue
                return images
            
            all_images = find_large_images(slide.shapes)
            slide_area = self.emu_to_px(self.prs.slide_width) * self.emu_to_px(self.prs.slide_height)
            
            # Сортируем по размеру (площади)
            all_images.sort(key=lambda x: x['area'], reverse=True)
            
            largest_image = None
            for img_data in all_images:
                area_percent = (img_data['area'] / slide_area) * 100
                
                # Если изображение занимает больше 40% слайда, вероятно это фон
                # Также проверяем, что оно не слишком далеко от начала слайда
                if area_percent > 40:
                    slide_width = self.emu_to_px(self.prs.slide_width)
                    slide_height = self.emu_to_px(self.prs.slide_height)
                    
                    # Позиция в процентах
                    left_percent = (img_data['left'] / slide_width) * 100
                    top_percent = (img_data['top'] / slide_height) * 100
                    
                    # Фон обычно начинается близко к краям (< 30% от левого/верхнего края)
                    if left_percent < 30 and top_percent < 30:
                        largest_image = img_data['shape']
                        print(f"  Обнаружен кандидат на фон: {area_percent:.1f}% слайда, позиция ({left_percent:.1f}%, {top_percent:.1f}%)")
                        break
            
            if largest_image:
                try:
                    img_ext = largest_image.image.ext
                    img_filename = f"slide{slide_num}_background{img_ext}"
                    img_path = os.path.join(self.images_dir, img_filename)
                    
                    with open(img_path, 'wb') as f:
                        f.write(largest_image.image.blob)
                    
                    print(f"  ✓ Фон найден (большое изображение): {img_filename}")
                    return f"images/{img_filename}"
                except Exception as e:
                    print(f"  Ошибка сохранения большого изображения: {e}")
            
            # Метод 2: Проверяем fill.type через XML
            if hasattr(slide.background, 'fill'):
                fill = slide.background.fill
                if hasattr(fill, 'type') and fill.type == 6:  # PICTURE
                    bg_element = slide.background._element
                    namespaces = {
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    }
                    
                    blip_elements = bg_element.xpath('.//a:blip[@r:embed]', namespaces=namespaces)
                    
                    if blip_elements:
                        rId = blip_elements[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        
                        if rId and rId in slide.part.rels:
                            image_part = slide.part.rels[rId].target_part
                            img_ext = image_part.ext
                            img_filename = f"slide{slide_num}_background{img_ext}"
                            img_path = os.path.join(self.images_dir, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(image_part.blob)
                            
                            print(f"  ✓ Фон найден (XML blipFill): {img_filename}")
                            return f"images/{img_filename}"
            
            # Метод 3: Ищем неиспользованные изображения в relationships
            used_images = set()
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        if hasattr(shape, 'image'):
                            # Используем размер blob как уникальный идентификатор
                            used_images.add(len(shape.image.blob))
                    except:
                        pass
            
            for rel_id, rel in slide.part.rels.items():
                if 'image' in rel.reltype.lower():
                    try:
                        image_blob = rel.target_part.blob
                        
                        # Если размер изображения не совпадает ни с одним в shapes
                        if len(image_blob) not in used_images:
                            img_ext = rel.target_part.ext
                            img_filename = f"slide{slide_num}_background{img_ext}"
                            img_path = os.path.join(self.images_dir, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(image_blob)
                            
                            print(f"  ✓ Фон найден (неиспользованное изображение): {img_filename}")
                            return f"images/{img_filename}"
                    except:
                        continue
                        
        except Exception as e:
            print(f"  Предупреждение: не удалось сохранить фон слайда {slide_num}: {e}")
        
        return None
    
    def save_master_background(self, slide, slide_num):
        """Извлекает фон из slide master (для слайдов с BACKGROUND fill type)
        Возвращает tuple: (background_color, background_image)
        """
        bg_color = None
        bg_image = None
        
        try:
            # Получаем slide layout и slide master
            slide_layout = slide.slide_layout
            slide_master = slide_layout.slide_master
            
            # Метод 1: Ищем большие FREEFORM фигуры с цветной заливкой и PICTURE в slide layout shapes
            # Важно: сначала проверяем FREEFORM (могут быть фоном под изображением)
            slide_area = self.emu_to_px(self.prs.slide_width) * self.emu_to_px(self.prs.slide_height)
            
            # Проход 1: Ищем FREEFORM с solid fill (может быть цветным фоном под изображением)
            for shape in slide_layout.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                        width = self.emu_to_px(shape.width)
                        height = self.emu_to_px(shape.height)
                        area = width * height
                        area_percent = (area / slide_area) * 100
                        
                        # Если FREEFORM занимает больше 90% слайда, это фон
                        if area_percent > 90:
                            left = self.emu_to_px(shape.left)
                            top = self.emu_to_px(shape.top)
                            
                            if left < 10 and top < 10:  # Начинается с начала слайда
                                # Проверяем заливку
                                if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID
                                    bg_color = self.rgb_to_hex(shape.fill.fore_color)
                                    if bg_color:
                                        print(f"  ✓ Цвет фона из slide layout (FREEFORM): {bg_color}")
                                        # Не возвращаем сразу, продолжаем искать изображение
                                        break
                except:
                    continue
            
            # Проход 2: Ищем PICTURE (может быть поверх цветного фона)
            for shape in slide_layout.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        width = self.emu_to_px(shape.width)
                        height = self.emu_to_px(shape.height)
                        area = width * height
                        area_percent = (area / slide_area) * 100
                        
                        # Если изображение занимает больше 30% слайда, вероятно это фон
                        if area_percent > 30:
                            img_ext = shape.image.ext
                            if not img_ext.startswith('.'):
                                img_ext = '.' + img_ext
                            img_filename = f"slide{slide_num}_layout_bg{img_ext}"
                            img_path = os.path.join(self.images_dir, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(shape.image.blob)
                            
                            print(f"  ✓ Фон из slide layout (изображение): {img_filename}")
                            bg_image = f"images/{img_filename}"
                            # Если есть и цвет и изображение, возвращаем оба
                            if bg_color and bg_image:
                                print(f"  ✓ Комбинированный фон: цвет {bg_color} + изображение")
                            return (bg_color, bg_image)
                except:
                    continue
            
            # Если нашли только цвет (без изображения), возвращаем его
            if bg_color:
                return (bg_color, None)
            
            # Метод 2: Проверяем background fill в slide layout
            if hasattr(slide_layout, 'background') and hasattr(slide_layout.background, 'fill'):
                fill = slide_layout.background.fill
                if hasattr(fill, 'type'):
                    if fill.type == 1:  # SOLID
                        try:
                            bg_color = self.rgb_to_hex(fill.fore_color)
                            print(f"  ✓ Цвет фона из slide layout: {bg_color}")
                            return (bg_color, None)
                        except:
                            pass
                    elif fill.type == 6:  # PICTURE
                        try:
                            bg_element = slide_layout.background._element
                            namespaces = {
                                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                            }
                            
                            blip_elements = bg_element.xpath('.//a:blip[@r:embed]', namespaces=namespaces)
                            
                            if blip_elements:
                                rId = blip_elements[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                
                                if rId and rId in slide_layout.part.rels:
                                    image_part = slide_layout.part.rels[rId].target_part
                                    img_ext = image_part.ext
                                    if not img_ext.startswith('.'):
                                        img_ext = '.' + img_ext
                                    img_filename = f"slide{slide_num}_layout_bg{img_ext}"
                                    img_path = os.path.join(self.images_dir, img_filename)
                                    
                                    with open(img_path, 'wb') as f:
                                        f.write(image_part.blob)
                                    
                                    print(f"  ✓ Фон из slide layout (fill): {img_filename}")
                                    return (None, f"images/{img_filename}")
                        except Exception as e:
                            pass
            
            # Метод 3: Проверяем изображения в slide master shapes
            for shape in slide_master.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        width = self.emu_to_px(shape.width)
                        height = self.emu_to_px(shape.height)
                        area = width * height
                        area_percent = (area / slide_area) * 100
                        
                        # Если изображение занимает больше 30% слайда
                        if area_percent > 30:
                            img_ext = shape.image.ext
                            if not img_ext.startswith('.'):
                                img_ext = '.' + img_ext
                            img_filename = f"slide{slide_num}_master_bg{img_ext}"
                            img_path = os.path.join(self.images_dir, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(shape.image.blob)
                            
                            print(f"  ✓ Фон из slide master (изображение): {img_filename}")
                            return (None, f"images/{img_filename}")
                except:
                    continue
            
            # Метод 4: Проверяем фон в slide master
            if hasattr(slide_master, 'background') and hasattr(slide_master.background, 'fill'):
                fill = slide_master.background.fill
                if hasattr(fill, 'type'):
                    if fill.type == 1:  # SOLID
                        try:
                            bg_color = self.rgb_to_hex(fill.fore_color)
                            print(f"  ✓ Цвет фона из slide master: {bg_color}")
                            return (bg_color, None)
                        except:
                            pass
                    elif fill.type == 6:  # PICTURE
                        try:
                            bg_element = slide_master.background._element
                            namespaces = {
                                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                            }
                            
                            blip_elements = bg_element.xpath('.//a:blip[@r:embed]', namespaces=namespaces)
                            
                            if blip_elements:
                                rId = blip_elements[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                
                                if rId and rId in slide_master.part.rels:
                                    image_part = slide_master.part.rels[rId].target_part
                                    img_ext = image_part.ext
                                    if not img_ext.startswith('.'):
                                        img_ext = '.' + img_ext
                                    img_filename = f"slide{slide_num}_master_bg{img_ext}"
                                    img_path = os.path.join(self.images_dir, img_filename)
                                    
                                    with open(img_path, 'wb') as f:
                                        f.write(image_part.blob)
                                    
                                    print(f"  ✓ Фон из slide master (fill): {img_filename}")
                                    return (None, f"images/{img_filename}")
                        except Exception as e:
                            pass
        
        except Exception as e:
            pass
        
        return (None, None)
    
    def process_text_frame(self, text_frame, shape_style, slide_width, slide_height):
        """Обрабатывает текстовый фрейм с адаптивными отступами"""
        html_content = []
        
        try:
            # Вертикальное выравнивание
            if text_frame.vertical_anchor:
                v_align = self.get_vertical_alignment(text_frame.vertical_anchor)
                shape_style['display'] = 'flex'
                shape_style['flex-direction'] = 'column'
                shape_style['justify-content'] = v_align
            
            # Отступы в процентах для адаптивности
            if text_frame.margin_left:
                margin_left_px = self.emu_to_px(text_frame.margin_left)
                margin_left_percent = (margin_left_px / slide_width) * 100
                shape_style['padding-left'] = f"{margin_left_percent:.2f}%"
            
            if text_frame.margin_right:
                margin_right_px = self.emu_to_px(text_frame.margin_right)
                margin_right_percent = (margin_right_px / slide_width) * 100
                shape_style['padding-right'] = f"{margin_right_percent:.2f}%"
            
            if text_frame.margin_top:
                margin_top_px = self.emu_to_px(text_frame.margin_top)
                margin_top_percent = (margin_top_px / slide_height) * 100
                shape_style['padding-top'] = f"{margin_top_percent:.2f}%"
            
            if text_frame.margin_bottom:
                margin_bottom_px = self.emu_to_px(text_frame.margin_bottom)
                margin_bottom_percent = (margin_bottom_px / slide_height) * 100
                shape_style['padding-bottom'] = f"{margin_bottom_percent:.2f}%"
            
            # Обработка параграфов
            for paragraph in text_frame.paragraphs:
                para_style = self.extract_paragraph_formatting(paragraph)
                para_html = []
                
                for run in paragraph.runs:
                    run_style = self.extract_text_formatting(run)
                    style_str = '; '.join([f"{k}: {v}" for k, v in run_style.items()])
                    
                    text = run.text.replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br>')
                    if style_str:
                        para_html.append(f'<span style="{style_str}">{text}</span>')
                    else:
                        para_html.append(text)
                
                para_style_str = '; '.join([f"{k}: {v}" for k, v in para_style.items()])
                para_content = ''.join(para_html)
                
                if para_content.strip():  # Только если есть контент
                    if para_style_str:
                        html_content.append(f'<p style="{para_style_str}">{para_content}</p>')
                    else:
                        html_content.append(f'<p>{para_content}</p>')
        except Exception as e:
            print(f"  Ошибка обработки текста: {e}")
        
        return '\n'.join(html_content)
    
    def process_slide(self, slide, slide_num):
        """Обрабатывает один слайд"""
        print(f"Обработка слайда {slide_num}...")
        
        slide_width = self.emu_to_px(self.prs.slide_width)
        slide_height = self.emu_to_px(self.prs.slide_height)
        
        shapes_data = []
        img_counter = 0
        shape_counter = 0  # Счетчик для z-index
        
        # Фон слайда
        background = None
        background_image = None
        
        # Сначала проверяем прямой фон слайда
        try:
            # Проверяем фоновый цвет
            if slide.background.fill.type == 1:  # SOLID
                background = self.rgb_to_hex(slide.background.fill.fore_color)
            # Проверяем фоновое изображение
            elif slide.background.fill.type == 6:  # PICTURE
                background_image = self.save_background_image(slide, slide_num)
        except Exception as e:
            pass
        
        # Дополнительная проверка: ищем большую FREEFORM фигуру, которая может быть фоном
        # Это для случаев, когда фон - это просто цветной прямоугольник
        # ПРИОРИТЕТ выше, чем slide master!
        if not background and not background_image:
            try:
                slide_area = slide_width * slide_height
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                            # Проверяем размер фигуры
                            shape_area = self.emu_to_px(shape.width) * self.emu_to_px(shape.height)
                            area_percent = (shape_area / slide_area) * 100
                            
                            # Если фигура занимает больше 95% слайда и начинается с (0,0)
                            if area_percent > 95:
                                left = self.emu_to_px(shape.left)
                                top = self.emu_to_px(shape.top)
                                
                                if left < 5 and top < 5:  # Почти в начале слайда
                                    # Проверяем заливку
                                    if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID
                                        background = self.rgb_to_hex(shape.fill.fore_color)
                                        print(f"  ✓ Фон найден как FREEFORM: {background}")
                                        break
                    except:
                        continue
            except:
                pass
        
        # Если фон все еще не найден, проверяем slide master (самый низкий приоритет)
        if not background and not background_image:
            try:
                if slide.background.fill.type == 5:  # BACKGROUND
                    master_bg_color, master_bg_image = self.save_master_background(slide, slide_num)
                    if master_bg_color:
                        background = master_bg_color
                    if master_bg_image:
                        background_image = master_bg_image
            except Exception as e:
                print(f"  Предупреждение: не удалось обработать фон слайда: {e}")
        
        # Сохраняем цвет фона для определения дефолтного цвета текста
        self.current_slide_bg_color = background if background else '#FFFFFF'
        
        # Обработка фигур (включая группы)
        def is_qr_code_group(group_shape):
            """Проверяет, является ли группа составным QR-кодом
            
            Args:
                group_shape: Группа для проверки
            
            Returns:
                bool: True если группа содержит маленькие изображения/фигуры, формирующие QR
            """
            if group_shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                return False
            
            # Критерии составного QR-кода:
            # 1. Группа содержит 10+ элементов (обычно QR состоит из множества квадратиков)
            # 2. Группа имеет небольшой размер (< 150px)
            # 3. Группа примерно квадратная
            
            try:
                width_px = group_shape.width // 9525
                height_px = group_shape.height // 9525
                num_shapes = len(group_shape.shapes)
                
                # Проверяем размер
                if width_px > 150 or height_px > 150:
                    return False
                
                # Проверяем форму (примерно квадрат)
                ratio = width_px / height_px if height_px > 0 else 0
                if not (0.7 < ratio < 1.3):
                    return False
                
                # Проверяем количество элементов
                if num_shapes < 10:
                    return False
                
                # Проверяем, что большинство элементов - это FREEFORM или PICTURE (части QR)
                freeform_count = sum(1 for s in group_shape.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM)
                picture_count = sum(1 for s in group_shape.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
                
                if (freeform_count + picture_count) / num_shapes < 0.8:
                    return False
                
                return True
            except:
                return False
        
        def process_qr_group_as_image(group_shape):
            """Обрабатывает группу как составной QR-код
            
            Создает единое изображение из группы, объединяя все её части
            """
            nonlocal img_counter, shape_counter
            
            try:
                # Увеличиваем счетчики
                img_counter += 1
                shape_counter += 1
                
                # Получаем границы группы
                group_width_px = group_shape.width // 9525
                group_height_px = group_shape.height // 9525
                group_left_px = group_shape.left // 9525
                group_top_px = group_shape.top // 9525
                
                # Сохраняем информацию о группе как QR-код
                width_percent = (group_width_px / slide_width) * 100
                height_percent = (group_height_px / slide_height) * 100
                left_percent = (group_left_px / slide_width) * 100
                top_percent = (group_top_px / slide_height) * 100
                
                # Собираем информацию о всех частях группы
                parts = []
                for sub_shape in group_shape.shapes:
                    try:
                        part_data = {
                            'type': sub_shape.shape_type,
                            'left': sub_shape.left // 9525,
                            'top': sub_shape.top // 9525,
                            'width': sub_shape.width // 9525,
                            'height': sub_shape.height // 9525,
                        }
                        
                        # Для FREEFORM - сохраняем цвет заливки
                        if sub_shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                            try:
                                if sub_shape.fill.type == MSO_FILL_TYPE.SOLID:
                                    rgb = sub_shape.fill.fore_color.rgb
                                    part_data['fill_color'] = f'rgb({rgb[0]}, {rgb[1]}, {rgb[2]})'
                                else:
                                    part_data['fill_color'] = 'transparent'
                            except:
                                part_data['fill_color'] = 'transparent'
                        
                        # Для PICTURE - сохраняем путь к изображению
                        elif sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = sub_shape.image
                                ext = image.ext
                                img_data = image.blob
                                img_name = f"slide{slide_num + 1}_qrpart{len(parts) + 1}.{ext}"
                                img_path = os.path.join(self.images_dir, img_name)
                                os.makedirs(os.path.dirname(img_path), exist_ok=True)
                                with open(img_path, 'wb') as f:
                                    f.write(img_data)
                                part_data['image_path'] = f"images/{img_name}"
                            except:
                                part_data['image_path'] = None
                        
                        parts.append(part_data)
                    except Exception as e:
                        print(f"    Предупреждение: не удалось обработать часть группы: {e}")
                        continue
                
                shape_data = {
                    'type': 'qr-group',
                    'style': {
                        'position': 'absolute',
                        'left': f"{left_percent:.3f}%",
                        'top': f"{top_percent:.3f}%",
                        'width': f"{width_percent:.3f}%",
                        'height': f"{height_percent:.3f}%",
                        'z-index': str(shape_counter),
                    },
                    'content': '',
                    'image_type': 'qr-code',
                    'is_composite': True,
                    'num_parts': len(parts),
                    'actual_size': (group_width_px, group_height_px),
                    'parts': parts,  # Список всех частей группы
                    'group_bounds': {  # Границы группы для расчета относительных позиций
                        'left': group_left_px,
                        'top': group_top_px,
                        'width': group_width_px,
                        'height': group_height_px
                    }
                }
                
                shapes_data.append(shape_data)
                print(f"  QR-группа: {group_width_px}x{group_height_px}px ({len(parts)} частей) → composite qr-code")
                print(f"    → Части: {len([p for p in parts if p['type'] == MSO_SHAPE_TYPE.FREEFORM])} FREEFORM, {len([p for p in parts if p['type'] == MSO_SHAPE_TYPE.PICTURE])} PICTURE")
                
            except Exception as e:
                print(f"  Предупреждение: не удалось обработать QR-группу: {e}")
        
        def process_shape_recursive(shape, level=0):
            """Рекурсивно обрабатывает фигуры, включая группы
            
            Args:
                shape: Фигура для обработки
                level: Уровень вложенности (0 = слайд, 1+ = внутри группы)
            
            Note:
                Координаты фигур внутри групп в python-pptx уже абсолютные относительно слайда,
                поэтому нет необходимости добавлять offset группы.
            """
            nonlocal img_counter, shape_counter
            
            # Пропускаем FREEFORM фигуры, которые использованы как фон слайда
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM and background and level == 0:
                try:
                    # Проверяем, это ли фоновая фигура (только на верхнем уровне слайда)
                    shape_area = self.emu_to_px(shape.width) * self.emu_to_px(shape.height)
                    slide_area = slide_width * slide_height
                    area_percent = (shape_area / slide_area) * 100
                    
                    if area_percent > 95:
                        left = self.emu_to_px(shape.left)
                        top = self.emu_to_px(shape.top)
                        
                        if left < 5 and top < 5:
                            # Это фоновая фигура, пропускаем
                            return
                except:
                    pass
            
            # Проверяем, является ли это группой
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Проверяем, является ли группа составным QR-кодом
                if is_qr_code_group(shape):
                    # Обрабатываем как единый QR-код
                    process_qr_group_as_image(shape)
                else:
                    # Обычная группа - обрабатываем каждую фигуру рекурсивно
                    # Координаты дочерних элементов уже абсолютные!
                    for sub_shape in shape.shapes:
                        process_shape_recursive(sub_shape, level + 1)
                return
            
            shape_data = {
                'type': None,
                'style': {},
                'content': ''
            }
            
            # Увеличиваем счетчик фигур для z-index
            shape_counter += 1
            current_shape_index = shape_counter
            
            try:
                # Координаты shape уже абсолютные, offset не нужен
                base_style = self.extract_shape_style(shape, slide_width, slide_height, current_shape_index)
            except Exception as e:
                print(f"  Предупреждение: не удалось извлечь стиль фигуры: {e}")
                return
            
            # Текстовые блоки
            if shape.has_text_frame and hasattr(shape, 'text') and shape.text.strip():
                shape_data['type'] = 'text'
                shape_data['style'] = base_style
                shape_data['content'] = self.process_text_frame(shape.text_frame, base_style, slide_width, slide_height)
                if shape_data['content'].strip():  # Только если есть контент
                    shapes_data.append(shape_data)
            
            # Изображения
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_counter += 1
                try:
                    img_path = self.save_image(shape.image, slide_num, img_counter, "img")
                    
                    if img_path:
                        # Для изображений создаём стиль БЕЗ background-color
                        # Копируем base_style и удаляем фон, чтобы сохранить прозрачность PNG
                        image_style = base_style.copy()
                        if 'background-color' in image_style:
                            del image_style['background-color']
                        if 'background' in image_style:
                            del image_style['background']
                        if 'opacity' in image_style:
                            # Opacity для картинок сохраняем
                            pass
                        
                        shape_data['type'] = 'image'
                        shape_data['style'] = image_style
                        shape_data['content'] = img_path
                        
                        # v15: Классификация изображения
                        try:
                            from PIL import Image
                            full_path = os.path.join(self.images_dir, os.path.basename(img_path))
                            
                            # Получаем позицию на слайде
                            left_percent = float(base_style['left'].rstrip('%'))
                            top_percent = float(base_style['top'].rstrip('%'))
                            width_px = shape.width // 9525
                            height_px = shape.height // 9525
                            
                            # Классифицируем изображение
                            classification = self.image_classifier.classify(
                                full_path,
                                (left_percent, top_percent),
                                (width_px, height_px)
                            )
                            
                            img_type = classification['type']
                            actual_w, actual_h = classification['actual_size']
                            
                            # Сохраняем классификацию
                            shape_data['image_type'] = img_type
                            shape_data['actual_size'] = (actual_w, actual_h)
                            shape_data['classification_confidence'] = classification['confidence']
                            
                            print(f"  Изображение: {actual_w}x{actual_h}px → {img_type} ({classification['confidence']:.0%})")
                            
                            # Для QR-кодов сохраняем флаг is_small для обратной совместимости
                            if img_type == 'qr-code':
                                shape_data['is_small'] = True
                                print(f"    → QR-код будет отображён в фактическом размере")
                            
                        except Exception as e_classify:
                            print(f"  Предупреждение: не удалось классифицировать изображение: {e_classify}")
                            # Fallback к старой логике
                            try:
                                with Image.open(full_path) as img:
                                    actual_w, actual_h = img.size
                                shape_data['actual_size'] = (actual_w, actual_h)
                                shape_data['image_type'] = 'unknown'
                            except:
                                pass
                        
                        shapes_data.append(shape_data)
                except Exception as e:
                    print(f"  Предупреждение: не удалось сохранить изображение: {e}")
            
            # Placeholder'ы (могут содержать стили даже если пустые)
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Проверяем, есть ли текст
                has_text = hasattr(shape, 'text') and shape.text.strip()
                
                # Проверяем, есть ли визуальные стили (фон, границы и т.д.)
                has_fill = hasattr(shape, 'fill') and shape.fill.type is not None
                has_line = hasattr(shape, 'line') and hasattr(shape.line, 'color')
                
                # Если есть текст
                if has_text:
                    shape_data['type'] = 'text'
                    shape_data['style'] = base_style
                    if shape.has_text_frame:
                        shape_data['content'] = self.process_text_frame(shape.text_frame, base_style, slide_width, slide_height)
                    else:
                        shape_data['content'] = f'<p>{shape.text}</p>'
                    shapes_data.append(shape_data)
                # Если нет текста, но есть стили (фон или граница)
                elif has_fill or has_line:
                    shape_data['type'] = 'shape'
                    shape_data['style'] = base_style
                    shape_data['content'] = ''
                    shapes_data.append(shape_data)
                    print(f"  Обработан пустой placeholder со стилями: {shape.name}")
            
            # Таблицы
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_data['type'] = 'table'
                shape_data['style'] = base_style
                shape_data['content'] = self.process_table(shape.table)
                shapes_data.append(shape_data)
            
            # Автофигуры с заливкой (прямоугольники, эллипсы и т.д.)
            elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, 
                                      MSO_SHAPE_TYPE.FREEFORM,
                                      MSO_SHAPE_TYPE.LINE,
                                      MSO_SHAPE_TYPE.TEXT_BOX]:
                # Если есть текст
                if hasattr(shape, 'text') and shape.text.strip():
                    shape_data['type'] = 'text'
                    shape_data['style'] = base_style
                    if shape.has_text_frame:
                        shape_data['content'] = self.process_text_frame(shape.text_frame, base_style, slide_width, slide_height)
                    else:
                        shape_data['content'] = f'<p>{shape.text}</p>'
                    shapes_data.append(shape_data)
                # Если нет текста, но есть заливка - фигура с фоном/границей
                else:
                    try:
                        # Проверяем, есть ли заливка-изображение
                        has_fill = hasattr(shape, 'fill') and shape.fill.type != None
                        
                        if has_fill:
                            try:
                                # Тип 6 = PICTURE (заливка изображением)
                                if shape.fill.type == 6:
                                    # Пытаемся извлечь изображение из заливки
                                    try:
                                        # Получаем blip (binary large image part) из заливки
                                        fill_element = shape.fill._fill
                                        if hasattr(fill_element, 'blipFill'):
                                            blip = fill_element.blipFill.blip
                                            if hasattr(blip, 'embed'):
                                                rId = blip.embed
                                                image_part = shape.part.related_part(rId)
                                                
                                                img_counter += 1
                                                ext = image_part.ext
                                                img_filename = f"slide{slide_num}_img{img_counter}{ext}"
                                                img_path = os.path.join(self.images_dir, img_filename)
                                                
                                                with open(img_path, 'wb') as f:
                                                    f.write(image_part.blob)
                                                
                                                print(f"  ✓ Сохранена заливка-изображение: {img_filename}")
                                                
                                                shape_data['type'] = 'image'
                                                shape_data['style'] = base_style
                                                shape_data['content'] = f"images/{img_filename}"
                                                shapes_data.append(shape_data)
                                                return  # Выходим, изображение обработано
                                    except Exception as e:
                                        print(f"  Предупреждение: не удалось извлечь заливку-изображение: {e}")
                                
                                # Если не изображение, обрабатываем как обычную фигуру
                                # base_style уже содержит все стили из extract_shape_style (background-color, border, opacity и т.д.)
                                if shape.fill.type == 1:  # SOLID
                                    shape_data['type'] = 'shape'
                                    shape_data['style'] = base_style  # Используем уже извлеченные стили
                                    shape_data['content'] = ''
                                    shapes_data.append(shape_data)
                            except Exception as e_fill:
                                pass
                        
                        # Если нет заливки, но есть граница или другие визуальные стили
                        # base_style уже содержит все стили из extract_shape_style
                        elif not has_fill:
                            has_line = hasattr(shape, 'line') and hasattr(shape.line, 'color') and shape.line.color
                            
                            if has_line:
                                shape_data['type'] = 'shape'
                                shape_data['style'] = base_style  # Используем уже извлеченные стили (включая border)
                                shape_data['content'] = ''
                                shapes_data.append(shape_data)
                    except:
                        pass  # Пропускаем фигуры, которые не можем обработать
            
            # Другие типы фигур
            else:
                if hasattr(shape, 'text') and shape.text.strip():
                    shape_data['type'] = 'shape'
                    shape_data['style'] = base_style
                    shape_data['content'] = shape.text
                    shapes_data.append(shape_data)
        
        # Обрабатываем все фигуры на слайде
        for shape in slide.shapes:
            process_shape_recursive(shape)
        
        return {
            'slide_num': slide_num,
            'width': slide_width,
            'height': slide_height,
            'aspect_ratio': slide_width / slide_height,
            'background': background,
            'background_image': background_image,
            'shapes': shapes_data
        }
    
    def process_table(self, table):
        """Обрабатывает таблицу"""
        html = ['<table style="width: 100%; border-collapse: collapse;">']
        
        try:
            for row in table.rows:
                html.append('<tr>')
                for cell in row.cells:
                    # Стили ячейки
                    cell_style = []
                    
                    if cell.fill.type == 1:  # SOLID
                        bg_color = self.rgb_to_hex(cell.fill.fore_color)
                        if bg_color:
                            cell_style.append(f"background-color: {bg_color}")
                    
                    # Границы
                    cell_style.append("border: 1px solid #ccc")
                    cell_style.append("padding: 8px")
                    
                    style_str = '; '.join(cell_style)
                    text = cell.text.replace('<', '&lt;').replace('>', '&gt;')
                    
                    html.append(f'<td style="{style_str}">{text}</td>')
                html.append('</tr>')
        except Exception as e:
            print(f"Ошибка обработки таблицы: {e}")
        
        html.append('</table>')
        return '\n'.join(html)
    
    def convert(self):
        """Основной метод конвертации"""
        self.load_presentation()
        
        # Обработка всех слайдов
        for idx, slide in enumerate(self.prs.slides, 1):
            slide_data = self.process_slide(slide, idx)
            self.slide_data.append(slide_data)
        
        # Генерация HTML
        self.generate_html()
        self.generate_css()
        
        # Сохранение метаданных
        self.save_metadata()
        
        print(f"\n✅ Конвертация завершена!")
        print(f"📁 Результаты сохранены в: {self.output_dir}")
        print(f"🌐 Откройте: {os.path.join(self.output_dir, 'index.html')}")
    
    def generate_html(self):
        """Генерирует HTML файлы - отдельный файл для каждого слайда"""
        
        # Генерируем отдельную страницу для каждого слайда
        for slide_data in self.slide_data:
            self._generate_slide_page(slide_data)
        
        # Генерируем главную страницу index.html со списком страниц
        self._generate_index_page()
        
        print(f"✅ Создано {len(self.slide_data)} HTML страниц в папке pages/")
    
    def _generate_slide_html_content(self, slide_data):
        """Генерирует HTML контент для одного слайда"""
        slide_num = slide_data['slide_num']
        aspect_ratio = slide_data['aspect_ratio']
        total_slides = len(self.slide_data)
        
        # Стили фона
        bg_styles = []
        if slide_data.get('background'):
            bg_styles.append(f"background-color: {slide_data['background']}")
        
        if slide_data.get('background_image'):
            # Путь к изображению должен быть относительно pages/
            bg_styles.append(f"background-image: url('../{slide_data['background_image']}')")
            bg_styles.append("background-size: cover")
            bg_styles.append("background-position: center")
            bg_styles.append("background-repeat: no-repeat")
        
        bg_style = '; '.join(bg_styles)
        
        html_parts = []
        
        # Фигуры на слайде
        for shape in slide_data['shapes']:
            style_str = '; '.join([f"{k}: {v}" for k, v in shape['style'].items()])
            
            if shape['type'] == 'text':
                html_parts.append(f'''
                <div class="text-block" style="{style_str}">
                    {shape['content']}
                </div>
''')
            elif shape['type'] == 'qr-group':
                # v16.3: Композитный QR-код из группы фигур
                parts = shape.get('parts', [])
                group_bounds = shape.get('group_bounds', {})
                
                html_parts.append(f'''
                <div class="qr-group-block" style="{style_str}; overflow: visible;">
''')
                
                for part in parts:
                    rel_left = ((part['left'] - group_bounds['left']) / group_bounds['width']) * 100
                    rel_top = ((part['top'] - group_bounds['top']) / group_bounds['height']) * 100
                    rel_width = (part['width'] / group_bounds['width']) * 100
                    rel_height = (part['height'] / group_bounds['height']) * 100
                    
                    part_style = f"position: absolute; left: {rel_left:.3f}%; top: {rel_top:.3f}%; width: {rel_width:.3f}%; height: {rel_height:.3f}%;"
                    
                    if part['type'] == MSO_SHAPE_TYPE.FREEFORM:
                        fill_color = part.get('fill_color', 'transparent')
                        html_parts.append(f'''
                    <div class="qr-part qr-freeform" style="{part_style} background-color: {fill_color};"></div>
''')
                    elif part['type'] == MSO_SHAPE_TYPE.PICTURE:
                        img_path = part.get('image_path')
                        if img_path:
                            # Корректируем путь для pages/
                            img_path = '../' + img_path
                            html_parts.append(f'''
                    <div class="qr-part qr-picture" style="{part_style}">
                        <img src="{img_path}" alt="QR Part" style="width: 100%; height: 100%; object-fit: contain; image-rendering: pixelated;">
                    </div>
''')
                
                html_parts.append('''
                </div>
''')
            elif shape['type'] == 'image':
                img_type = shape.get('image_type', 'unknown')
                actual_w, actual_h = shape.get('actual_size', (0, 0))
                # Корректируем путь к изображению для pages/
                img_src = '../' + shape['content']
                
                if img_type == 'qr-code':
                    html_parts.append(f'''
                <div class="image-block qr-code" style="{style_str}; display: flex; align-items: center; justify-content: center;">
                    <img src="{img_src}" alt="QR Code" style="width: {actual_w}px; height: {actual_h}px; object-fit: none; image-rendering: pixelated;">
                </div>
''')
                elif img_type == 'icon':
                    html_parts.append(f'''
                <div class="image-block icon" style="{style_str}; display: flex; align-items: center; justify-content: center;">
                    <img src="{img_src}" alt="Icon" style="max-width: 100%; max-height: 100%; object-fit: contain;">
                </div>
''')
                elif img_type == 'logo':
                    html_parts.append(f'''
                <div class="image-block logo" style="{style_str}">
                    <img src="{img_src}" alt="Logo" style="width: 100%; height: 100%; object-fit: contain;">
                </div>
''')
                elif img_type == 'diagram':
                    html_parts.append(f'''
                <div class="image-block diagram" style="{style_str}">
                    <img src="{img_src}" alt="Diagram" style="width: 100%; height: 100%; object-fit: contain;">
                </div>
''')
                else:
                    if shape.get('is_small', False) and actual_w > 0:
                        html_parts.append(f'''
                <div class="image-block" style="{style_str}; display: flex; align-items: center; justify-content: center;">
                    <img src="{img_src}" alt="Image" style="width: {actual_w}px; height: {actual_h}px; object-fit: none;">
                </div>
''')
                    else:
                        html_parts.append(f'''
                <div class="image-block" style="{style_str}">
                    <img src="{img_src}" alt="Image" style="width: 100%; height: 100%; object-fit: contain;">
                </div>
''')
            elif shape['type'] == 'table':
                html_parts.append(f'''
                <div class="table-block" style="{style_str}">
                    {shape['content']}
                </div>
''')
            elif shape['type'] == 'group':
                html_parts.append(f'''
                <div class="group-block" style="{style_str}">
''')
                for sub_shape in shape['content']:
                    sub_style_str = '; '.join([f"{k}: {v}" for k, v in sub_shape['style'].items()])
                    
                    if sub_shape['type'] == 'shape':
                        html_parts.append(f'''
                    <div class="shape-block" style="{sub_style_str}"></div>
''')
                    elif sub_shape['type'] == 'image':
                        sub_img_src = '../' + sub_shape['content']
                        if sub_shape.get('is_small', False) and 'actual_size' in sub_shape:
                            actual_w, actual_h = sub_shape['actual_size']
                            html_parts.append(f'''
                    <div class="image-block" style="{sub_style_str}; display: flex; align-items: center; justify-content: center;">
                        <img src="{sub_img_src}" alt="Image" style="width: {actual_w}px; height: {actual_h}px; object-fit: none;">
                    </div>
''')
                        else:
                            html_parts.append(f'''
                    <div class="image-block" style="{sub_style_str}">
                        <img src="{sub_img_src}" alt="Image" style="width: 100%; height: 100%; object-fit: contain;">
                    </div>
''')
                    elif sub_shape['type'] == 'text':
                        html_parts.append(f'''
                    <div class="text-block" style="{sub_style_str}">
                        {sub_shape['content']}
                    </div>
''')
                
                html_parts.append('''                </div>
''')
            elif shape['type'] == 'shape':
                html_parts.append(f'''
                <div class="shape-block" style="{style_str}">
                    <p>{shape['content']}</p>
                </div>
''')
        
        return ''.join(html_parts), bg_style, aspect_ratio
    
    def _generate_slide_page(self, slide_data):
        """Генерирует HTML файл для одного слайда"""
        slide_num = slide_data['slide_num']
        total_slides = len(self.slide_data)
        slide_width = slide_data['width']
        slide_height = slide_data['height']
        
        # Получаем контент слайда
        slide_content, bg_style, aspect_ratio = self._generate_slide_html_content(slide_data)
        
        # Навигация к соседним слайдам
        prev_link = f'page{slide_num-1}.html' if slide_num > 1 else ''
        next_link = f'page{slide_num+1}.html' if slide_num < total_slides else ''
        
        html = f'''<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Страница {slide_num}</title>
    <link rel="stylesheet" href="../style.css">
    <style>
        /* Точные размеры слайда для этой страницы */
        .slide {{
            width: {slide_width}px;
            height: {slide_height}px;
            max-width: 95vw;
            max-height: 85vh;
        }}
        
        /* Адаптивное масштабирование при необходимости */
        @media (max-width: {slide_width}px), (max-height: {slide_height}px) {{
            .slide {{
                width: 95vw;
                height: calc(95vw * {aspect_ratio:.6f});
                max-height: 85vh;
            }}
            
            @supports (width: min(95vw, {slide_width}px)) {{
                .slide {{
                    width: min(95vw, {slide_width}px);
                    height: min(calc(95vw * {aspect_ratio:.6f}), {slide_height}px);
                    max-height: 85vh;
                }}
            }}
        }}
    </style>
</head>
<body>
    <div class="presentation-container">
        <!-- Navigation -->
        <nav class="presentation-nav">
            <a href="../index.html" class="nav-btn">📋 К списку</a>
            {f'<a href="{prev_link}" class="nav-btn">← Назад</a>' if prev_link else '<span class="nav-btn disabled">← Назад</span>'}
            <span class="slide-counter">{slide_num} / {total_slides}</span>
            {f'<a href="{next_link}" class="nav-btn">Вперед →</a>' if next_link else '<span class="nav-btn disabled">Вперед →</span>'}
        </nav>
        
        <!-- Slide Content -->
        <div class="slides-wrapper">
            <div class="slide active" data-slide="{slide_num}" data-aspect="{aspect_ratio:.4f}" data-width="{slide_width}" data-height="{slide_height}" style="{bg_style}">
{slide_content}
            </div>
        </div>
    </div>
    
    <script>
        // Keyboard navigation
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowLeft' && '{prev_link}') {{
                window.location.href = '{prev_link}';
            }} else if (e.key === 'ArrowRight' && '{next_link}') {{
                window.location.href = '{next_link}';
            }} else if (e.key === 'Escape') {{
                window.location.href = '../index.html';
            }}
        }});
        
        // Fullscreen toggle
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'F11') {{
                e.preventDefault();
                if (!document.fullscreenElement) {{
                    document.documentElement.requestFullscreen();
                }} else {{
                    document.exitFullscreen();
                }}
            }}
        }});
    </script>
</body>
</html>
'''
        
        # Сохраняем файл
        page_path = os.path.join(self.pages_dir, f'page{slide_num}.html')
        with open(page_path, 'w', encoding='utf-8') as f:
            f.write(html)
    
    def _generate_index_page(self):
        """Генерирует главную страницу index.html со списком всех слайдов"""
        html_parts = ['''<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Презентация - Список страниц</title>
    <link rel="stylesheet" href="style.css">
    <style>
        .page-list-container {
            max-width: 1200px;
            margin: 50px auto;
            padding: 20px;
        }
        
        .page-list-title {
            font-size: 32px;
            font-weight: bold;
            color: white;
            text-align: center;
            margin-bottom: 30px;
        }
        
        .page-grid {
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(250px, 1fr));
            gap: 20px;
            padding: 20px;
        }
        
        .page-card {
            background: white;
            border-radius: 10px;
            padding: 20px;
            text-align: center;
            transition: all 0.3s ease;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.3);
            cursor: pointer;
        }
        
        .page-card:hover {
            transform: translateY(-5px);
            box-shadow: 0 8px 12px rgba(76, 175, 80, 0.5);
        }
        
        .page-card a {
            text-decoration: none;
            color: #333;
            font-size: 20px;
            font-weight: 600;
            display: block;
        }
        
        .page-card .page-number {
            font-size: 48px;
            color: #4CAF50;
            margin-bottom: 10px;
        }
    </style>
</head>
<body>
    <div class="page-list-container">
        <h1 class="page-list-title">📚 Список страниц презентации</h1>
        <div class="page-grid">
''']
        
        for slide_data in self.slide_data:
            slide_num = slide_data['slide_num']
            html_parts.append(f'''
            <div class="page-card" onclick="window.location.href='pages/page{slide_num}.html'">
                <div class="page-number">{slide_num}</div>
                <a href="pages/page{slide_num}.html">Страница {slide_num}</a>
            </div>
''')
        
        html_parts.append('''
        </div>
    </div>
    
    <script>
        // Keyboard navigation - numbers 1-9 and 0
        document.addEventListener('keydown', (e) => {
            const key = e.key;
            if (key >= '1' && key <= '9') {
                const pageNum = parseInt(key);
                if (pageNum <= ''' + str(len(self.slide_data)) + ''') {
                    window.location.href = `pages/page${pageNum}.html`;
                }
            }
        });
    </script>
</body>
</html>
''')
        
        index_path = os.path.join(self.output_dir, 'index.html')
        with open(index_path, 'w', encoding='utf-8') as f:
            f.write(''.join(html_parts))
        
        print(f"✅ Главная страница создана: {index_path}")
    
    def generate_css(self):
        """Генерирует CSS файл"""
        
        css_content = '''/* PPTX to HTML - Generated Styles with Fixed Layout */

* {
    margin: 0;
    padding: 0;
    box-sizing: border-box;
}

body {
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    background-color: #1a1a1a;
    color: #333;
    overflow: hidden;
}

.presentation-container {
    position: relative;
    width: 100vw;
    height: 100vh;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
}

/* Navigation */
.presentation-nav {
    position: fixed;
    top: 20px;
    left: 50%;
    transform: translateX(-50%);
    display: flex;
    gap: 20px;
    align-items: center;
    background: rgba(0, 0, 0, 0.8);
    padding: 15px 30px;
    border-radius: 50px;
    z-index: 1000;
    backdrop-filter: blur(10px);
}

.nav-btn {
    background: #4CAF50;
    color: white;
    border: none;
    padding: 10px 20px;
    border-radius: 25px;
    cursor: pointer;
    font-size: 16px;
    font-weight: 600;
    transition: all 0.3s ease;
    text-decoration: none;
    display: inline-block;
}

.nav-btn:hover {
    background: #45a049;
    transform: scale(1.05);
}

.nav-btn:active {
    transform: scale(0.95);
}

.nav-btn.disabled {
    background: #666;
    cursor: not-allowed;
    opacity: 0.5;
}

.nav-btn.disabled:hover {
    background: #666;
    transform: none;
}

.slide-counter {
    color: white;
    font-size: 18px;
    font-weight: 600;
    min-width: 80px;
    text-align: center;
}

/* Slides - точные размеры задаются inline в каждом HTML файле */
.slides-wrapper {
    position: relative;
    display: flex;
    align-items: center;
    justify-content: center;
    width: 100%;
    height: 100%;
    padding: 100px 20px 20px;
}

.slide {
    position: relative;
    display: block;
    opacity: 1;
    background: white;
    box-shadow: 0 10px 50px rgba(0, 0, 0, 0.5);
    transform-origin: center center;
}

.slide.active {
    display: block;
    opacity: 1;
    animation: slideIn 0.5s ease;
}

@keyframes slideIn {
    from {
        opacity: 0;
        transform: scale(0.95);
    }
    to {
        opacity: 1;
        transform: scale(1);
    }
}

/* Text blocks - используют точные размеры из PPTX */
.text-block {
    overflow: hidden;
    box-sizing: border-box;
}

.text-block p {
    margin: 0;
    padding: 2px 0;
    word-wrap: break-word;
    overflow-wrap: break-word;
}

.text-block span {
    white-space: pre-wrap;
}

/* Image blocks - сохраняют пропорции */
.image-block {
    overflow: hidden;
}

.image-block img {
    display: block;
    width: 100% !important;
    height: 100% !important;
    object-fit: contain;
}

/* Table blocks */
.table-block {
    overflow: auto;
}

.table-block table {
    width: 100%;
    height: 100%;
}

/* Shape blocks */
.shape-block {
    display: flex;
    align-items: center;
    justify-content: center;
    text-align: center;
    word-wrap: break-word;
}

/* Thumbnails Panel */
.thumbnails-panel {
    position: fixed;
    right: -300px;
    top: 0;
    bottom: 0;
    width: 300px;
    background: rgba(0, 0, 0, 0.9);
    backdrop-filter: blur(10px);
    transition: right 0.3s ease;
    z-index: 999;
    display: flex;
    flex-direction: column;
    padding: 80px 20px 20px;
}

.thumbnails-panel.open {
    right: 0;
}

.thumbnails-toggle {
    position: absolute;
    left: -50px;
    top: 50%;
    transform: translateY(-50%);
    background: rgba(0, 0, 0, 0.8);
    color: white;
    border: none;
    padding: 15px;
    border-radius: 10px 0 0 10px;
    cursor: pointer;
    font-size: 24px;
    transition: all 0.3s ease;
}

.thumbnails-toggle:hover {
    background: rgba(0, 0, 0, 1);
    left: -55px;
}

.thumbnails-grid {
    display: flex;
    flex-direction: column;
    gap: 15px;
    overflow-y: auto;
    padding-right: 10px;
}

.thumbnails-grid::-webkit-scrollbar {
    width: 8px;
}

.thumbnails-grid::-webkit-scrollbar-track {
    background: rgba(255, 255, 255, 0.1);
    border-radius: 4px;
}

.thumbnails-grid::-webkit-scrollbar-thumb {
    background: rgba(255, 255, 255, 0.3);
    border-radius: 4px;
}

.thumbnails-grid::-webkit-scrollbar-thumb:hover {
    background: rgba(255, 255, 255, 0.5);
}

.thumbnail {
    background: white;
    padding: 10px;
    border-radius: 8px;
    cursor: pointer;
    transition: all 0.3s ease;
    text-align: center;
    font-weight: 600;
    border: 3px solid transparent;
}

.thumbnail:hover {
    transform: scale(1.05);
    border-color: #4CAF50;
}

.thumbnail.active {
    border-color: #4CAF50;
    box-shadow: 0 0 20px rgba(76, 175, 80, 0.5);
}

/* Responsive Design - только для очень маленьких экранов */
@media (max-width: 768px) {
    .slides-wrapper {
        padding: 80px 10px 10px;
    }
    
    .presentation-nav {
        padding: 10px 15px;
        gap: 10px;
    }
    
    .nav-btn {
        padding: 8px 15px;
        font-size: 14px;
    }
    
    .thumbnails-panel {
        width: 200px;
        right: -200px;
    }
}

@media (max-width: 480px) {
    .presentation-nav {
        flex-direction: column;
        gap: 5px;
        padding: 10px;
    }
}

/* Print Styles */
@media print {
    body {
        background: white;
    }
    
    .presentation-nav,
    .thumbnails-panel {
        display: none;
    }
    
    .slide {
        display: block !important;
        opacity: 1 !important;
        position: relative !important;
        page-break-after: always;
        box-shadow: none;
        margin: 20px auto;
    }
}

/* Accessibility */
.nav-btn:focus,
.thumbnails-toggle:focus,
.thumbnail:focus {
    outline: 3px solid #4CAF50;
    outline-offset: 2px;
}

/* Loading Animation */
@keyframes fadeIn {
    from {
        opacity: 0;
    }
    to {
        opacity: 1;
    }
}
'''
        
        css_path = os.path.join(self.output_dir, 'style.css')
        with open(css_path, 'w', encoding='utf-8') as f:
            f.write(css_content)
        
        print(f"✅ CSS создан: {css_path}")
    
    def save_metadata(self):
        """Сохраняет метаданные презентации"""
        metadata = {
            'source_file': self.pptx_path,
            'total_slides': len(self.slide_data),
            'slides': []
        }
        
        for slide in self.slide_data:
            slide_num = slide['slide_num']
            
            # Определяем пути к ресурсам слайда
            slide_meta = {
                'slide_num': slide_num,
                'width': slide['width'],
                'height': slide['height'],
                'shapes_count': len(slide['shapes']),
                'html_page': f'pages/page{slide_num}.html',  # Путь к отдельной странице
                'html_url': f'pages/page{slide_num}.html'  # Полный URL к странице
            }
            
            # Добавляем информацию о фоновом изображении, если есть
            if slide.get('background_image'):
                slide_meta['background_image'] = slide['background_image']
            
            # Добавляем информацию о фоновом цвете, если есть
            if slide.get('background_color'):
                slide_meta['background_color'] = slide['background_color']
            
            metadata['slides'].append(slide_meta)
        
        metadata_path = os.path.join(self.output_dir, 'metadata.json')
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)
        
        print(f"✅ Метаданные сохранены: {metadata_path}")


def main():
    """Главная функция"""
    import sys
    import io
    
    # Устанавливаем UTF-8 для вывода
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    
    print("=" * 60)
    print("PPTX to HTML Converter v17.0")
    print("Конвертер презентаций PowerPoint в веб-страницы")
    print("Каждый слайд сохраняется в отдельный HTML файл")
    print("=" * 60)
    print()
    
    # Парсим аргументы командной строки
    pptx_file = None
    output_dir = None
    args = sys.argv[1:]
    
    # Получаем путь к файлу
    if len(args) > 0:
        pptx_file = args[0]
    else:
        pptx_file = input("Введите путь к PPTX файлу: ").strip().strip('"')
    
    if not os.path.exists(pptx_file):
        print(f"❌ Файл не найден: {pptx_file}")
        return
    
    if not pptx_file.lower().endswith('.pptx'):
        print("❌ Файл должен иметь расширение .pptx")
        return
    
    # Получаем папку вывода
    if len(args) > 1:
        output_dir = args[1]
    else:
        output_dir = input("Папка для сохранения (Enter = 'pptx_output'): ").strip()
        if not output_dir:
            output_dir = 'pptx_output'
    
    print()
    print("🚀 Начинаем конвертацию...")
    print()
    
    try:
        converter = PPTXToHTMLConverter(pptx_file, output_dir)
        converter.convert()
        
        print()
        print("=" * 60)
        print("✨ Готово! Презентация успешно конвертирована!")
        print("=" * 60)
        print()
        print("📝 Инструкции:")
        print(f"   1. Откройте: {os.path.join(output_dir, 'index.html')}")
        print("   2. Выберите страницу из списка или используйте навигацию")
        print("   3. Каждая страница доступна по отдельной ссылке")
        print()
        
    except Exception as e:
        print(f"❌ Ошибка конвертации: {e}")
        import traceback
        traceback.print_exc()


if __name__ == '__main__':
    main()
