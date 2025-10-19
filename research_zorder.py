#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Исследование z-order в PPTX
Проверяем, как PowerPoint определяет порядок наложения фигур
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_shape_order(shape, index, slide_width, slide_height, level=0):
    """Анализирует порядок фигуры в XML"""
    indent = "  " * level
    
    left_pct = (shape.left / slide_width) * 100
    top_pct = (shape.top / slide_height) * 100
    
    shape_type = MSO_SHAPE_TYPE(shape.shape_type).name
    
    # Пытаемся получить z-order из XML
    z_order = None
    try:
        # Проверяем атрибут в XML
        elem = shape._element
        # Ищем атрибут order или аналогичный
        if hasattr(elem, 'attrib'):
            z_order = elem.attrib.get('order', None)
        
        # Проверяем другие возможные источники z-order
        if not z_order and hasattr(elem, 'get'):
            z_order = elem.get('order', None)
    except:
        pass
    
    print(f"{indent}#{index}: {shape.name} ({shape_type})")
    print(f"{indent}   Position: left={left_pct:.1f}%, top={top_pct:.1f}%")
    print(f"{indent}   Order in shapes: {index}")
    print(f"{indent}   Z-order from XML: {z_order if z_order else 'N/A'}")
    
    # Проверяем текст
    has_text = False
    if hasattr(shape, 'text'):
        try:
            if shape.text.strip():
                has_text = True
                print(f"{indent}   Has text: YES ({shape.text[:30]}...)")
        except:
            pass
    
    if not has_text:
        print(f"{indent}   Has text: NO")
    
    return {
        'index': index,
        'name': shape.name,
        'type': shape_type,
        'left': left_pct,
        'top': top_pct,
        'xml_order': z_order,
        'has_text': has_text
    }

# Загружаем презентацию
prs = Presentation('NEO INVESTMENTS-fin.pptx')

print("=" * 70)
print("ИССЛЕДОВАНИЕ Z-ORDER В PPTX")
print("=" * 70)

slide_width = prs.slide_width
slide_height = prs.slide_height

for slide_num in [8, 10]:  # Проверяем слайды 8 и 10
    slide = prs.slides[slide_num - 1]
    
    print(f"\n{'='*70}")
    print(f"СЛАЙД #{slide_num}")
    print(f"{'='*70}")
    print(f"Всего фигур: {len(slide.shapes)}")
    
    all_shapes = []
    
    for i, shape in enumerate(slide.shapes, 1):
        print()
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result = analyze_shape_order(shape, i, slide_width, slide_height, 0)
            all_shapes.append(result)
            
            # Анализируем содержимое группы
            print(f"   GROUP contains {len(shape.shapes)} shapes:")
            for j, sub_shape in enumerate(shape.shapes, 1):
                print()
                sub_result = analyze_shape_order(sub_shape, f"{i}.{j}", slide_width, slide_height, 1)
                all_shapes.append(sub_result)
        else:
            result = analyze_shape_order(shape, i, slide_width, slide_height, 0)
            all_shapes.append(result)
    
    print(f"\n{'='*70}")
    print(f"ВЫВОД ДЛЯ СЛАЙДА {slide_num}:")
    print(f"{'='*70}")
    print("\n💡 ВАЖНО: Порядок фигур в slide.shapes определяет z-order!")
    print("   - Фигуры в начале списка (индекс 1, 2, 3...) на ЗАДНЕМ плане")
    print("   - Фигуры в конце списка (последние) на ПЕРЕДНЕМ плане")
    print("\n📊 Рекомендация: Использовать индекс в slide.shapes как базу для z-index")

print("\n" + "=" * 70)
print("ОБЩИЙ ВЫВОД:")
print("=" * 70)
print("""
В PowerPoint порядок наложения фигур определяется их позицией в коллекции shapes:
1. Первая фигура в списке (shapes[0]) - самая задняя (z-index = 1)
2. Последняя фигура в списке (shapes[-1]) - самая передняя (z-index = max)

Для групп:
- Группа имеет свой z-index относительно других фигур на слайде
- Внутри группы фигуры также упорядочены (первая = задняя, последняя = передняя)

ПРАВИЛЬНАЯ ФОРМУЛА Z-INDEX:
z-index = shape_index (порядок обработки в slide.shapes)

Дополнительно можно учесть вертикальную позицию для естественности:
z-index = shape_index * 10 + int(top_percent)
""")
