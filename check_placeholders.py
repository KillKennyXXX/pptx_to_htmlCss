"""
Проверка содержимого placeholder'ов на слайде 1
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_px(emu):
    return emu / 9525

# Загружаем презентацию
prs = Presentation("NEO INVESTMENTS-fin.pptx")
slide = prs.slides[0]

print("=" * 70)
print("ПРОВЕРКА PLACEHOLDER'ОВ НА СЛАЙДЕ 1")
print("=" * 70)
print()

# Проверяем фигуры на самом слайде (не layout)
print("Фигуры в slide.shapes:")
for idx, shape in enumerate(slide.shapes, 1):
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        print(f"\nPlaceholder #{idx}:")
        print(f"  Name: {shape.name}")
        
        # Проверяем текст
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            print(f"  Text: '{text}' (пустой={len(text)==0})")
        elif hasattr(shape, 'text'):
            text = shape.text
            print(f"  Text: '{text}' (пустой={len(text)==0})")
        else:
            print(f"  Text: НЕТ ТЕКСТА")
        
        # Позиция
        top_percent = (emu_to_px(shape.top) / emu_to_px(prs.slide_height)) * 100
        print(f"  Position: {top_percent:.1f}% от верха")

print("\n" + "=" * 70)
print("ВЫВОД:")
print("Если placeholder'ы пустые, они не должны отображаться")
print("Если имеют содержимое, их нужно обработать")
print("=" * 70)
