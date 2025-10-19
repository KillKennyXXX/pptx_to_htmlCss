"""
Тест координат фигур в группе
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation("NEO INVESTMENTS-fin.pptx")
slide = prs.slides[10]  # Слайд 11 (индекс 10)

print("="*80)
print("АНАЛИЗ КООРДИНАТ ГРУППЫ И ЕЁ ДОЧЕРНИХ ЭЛЕМЕНТОВ")
print("="*80)

for shape in slide.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"\n📦 ГРУППА: {shape.name}")
        print(f"   Позиция группы: left={shape.left}, top={shape.top}")
        print(f"   Размер группы: width={shape.width}, height={shape.height}")
        print(f"\n   Дочерние элементы:")
        
        for i, sub_shape in enumerate(shape.shapes, 1):
            print(f"\n   [{i}] {sub_shape.name}")
            print(f"       left={sub_shape.left}, top={sub_shape.top}")
            print(f"       width={sub_shape.width}, height={sub_shape.height}")
            
            # Проверяем тип
            shape_type_name = str(sub_shape.shape_type).split('.')[-1].split(' ')[0]
            print(f"       тип: {shape_type_name}")
            
            # Если это FREEFORM с заливкой
            if sub_shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                if hasattr(sub_shape, 'fill') and sub_shape.fill.type == 1:
                    try:
                        r, g, b = sub_shape.fill.fore_color.rgb
                        hex_color = f'#{r:02x}{g:02x}{b:02x}'
                        print(f"       цвет: RGB({r}, {g}, {b}) = {hex_color}")
                    except:
                        print(f"       цвет: не удалось получить")

print("\n" + "="*80)
